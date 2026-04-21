#!/usr/bin/env python3
"""
html_to_pptx.py — Convert 1920×1080 HTML slides into an editable PPTX.

Each HTML file becomes one slide. Native PowerPoint text boxes + rectangles
are emitted, so text stays editable in PowerPoint / Keynote.

Usage:
    python html_to_pptx.py "slides/*.html" -o deck.pptx
    python html_to_pptx.py slide1.html slide2.html -o deck.pptx
    python html_to_pptx.py "slides/*.html" -o deck.pptx --width 1920 --height 1080

Requirements:
    pip install playwright python-pptx
    playwright install chromium
"""

from __future__ import annotations

import argparse
import asyncio
import glob
import json
import re
import sys
from dataclasses import dataclass
from pathlib import Path

from playwright.async_api import async_playwright, Page, Browser

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

# ----------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------

EMU_PER_PX = 9525  # 1 px ≈ 9525 EMU at 96 DPI


def px_to_emu(px: float) -> int:
    return int(round(px * EMU_PER_PX))


def parse_rgb(s: str) -> tuple[int, int, int, float] | None:
    """Parse 'rgb(r,g,b)' / 'rgba(r,g,b,a)' / '#rrggbb' / '#rgb' → (r, g, b, alpha)."""
    if not s:
        return None
    s = s.strip().lower()
    if s in ("transparent", "none", "inherit", "currentcolor"):
        return None
    m = re.match(r"rgba?\(\s*([\d.]+)\s*,\s*([\d.]+)\s*,\s*([\d.]+)\s*(?:,\s*([\d.]+))?\s*\)", s)
    if m:
        r, g, b = int(float(m.group(1))), int(float(m.group(2))), int(float(m.group(3)))
        a = float(m.group(4)) if m.group(4) is not None else 1.0
        return (r, g, b, a)
    m = re.match(r"#([0-9a-f]{3})$", s)
    if m:
        r, g, b = (int(c * 2, 16) for c in m.group(1))
        return (r, g, b, 1.0)
    m = re.match(r"#([0-9a-f]{6})$", s)
    if m:
        r = int(m.group(1)[0:2], 16)
        g = int(m.group(1)[2:4], 16)
        b = int(m.group(1)[4:6], 16)
        return (r, g, b, 1.0)
    return None


def rgb_to_color(rgb: tuple[int, int, int, float]) -> RGBColor:
    return RGBColor(rgb[0], rgb[1], rgb[2])


def first_font(font_family: str) -> str:
    """'IBM Plex Sans', Helvetica, Arial, sans-serif → IBM Plex Sans"""
    if not font_family:
        return "Arial"
    first = font_family.split(",")[0].strip().strip('"').strip("'")
    return first or "Arial"


def css_align_to_pp(v: str) -> int:
    return {
        "left": PP_ALIGN.LEFT,
        "right": PP_ALIGN.RIGHT,
        "center": PP_ALIGN.CENTER,
        "justify": PP_ALIGN.JUSTIFY,
        "start": PP_ALIGN.LEFT,
        "end": PP_ALIGN.RIGHT,
    }.get((v or "").lower(), PP_ALIGN.LEFT)


# ----------------------------------------------------------------------
# DOM extraction (runs inside the page via page.evaluate)
# ----------------------------------------------------------------------

EXTRACT_JS = r"""
() => {
  // --- speaker notes --------------------------------------------------
  const notesEl = document.querySelector('script[type="application/json"]#speaker-notes');
  let speakerNotes = "";
  if (notesEl) {
    try {
      const arr = JSON.parse(notesEl.textContent);
      if (Array.isArray(arr)) speakerNotes = arr[0] || "";
      else if (typeof arr === "string") speakerNotes = arr;
    } catch (e) {}
  }

  // --- find the slide canvas & unscale it ----------------------------
  const canvas =
    document.querySelector('.canvas') ||
    document.querySelector('.slide') ||
    document.body;
  // Remove any transform so coords are 1:1
  const stage = document.querySelector('.stage');
  if (stage) {
    stage.style.position = 'static';
    stage.style.width = canvas.offsetWidth + 'px';
    stage.style.height = canvas.offsetHeight + 'px';
    stage.style.inset = 'auto';
    stage.style.background = '#fff';
  }
  canvas.style.transform = 'none';
  canvas.style.transformOrigin = 'top left';
  canvas.style.margin = '0';
  canvas.style.position = 'static';

  // Also remove the stage fixed positioning that causes issues
  document.querySelectorAll('.stage').forEach(s => {
    s.style.position = 'static';
    s.style.display = 'block';
    s.style.overflow = 'visible';
  });

  // Hide tweaks panel and other non-slide elements
  document.querySelectorAll('#tweaks, .export-hidden').forEach(el => {
    el.style.display = 'none';
  });

  const canvasRect = canvas.getBoundingClientRect();
  const originX = canvasRect.left;
  const originY = canvasRect.top;

  const shapes = [];   // background rects / borders
  const texts = [];    // text runs

  function visible(el, rect) {
    if (!rect || rect.width < 1 || rect.height < 1) return false;
    const s = getComputedStyle(el);
    if (s.display === 'none' || s.visibility === 'hidden' || parseFloat(s.opacity) === 0) return false;
    return true;
  }

  function walkShapes(el) {
    if (!(el instanceof Element)) return;
    if (el.id === 'tweaks' || el.tagName === 'SCRIPT' || el.tagName === 'STYLE') return;
    const s = getComputedStyle(el);
    const rect = el.getBoundingClientRect();
    if (!visible(el, rect)) return;

    const bg = s.backgroundColor;
    const border = {
      t: { w: parseFloat(s.borderTopWidth) || 0, c: s.borderTopColor, st: s.borderTopStyle },
      r: { w: parseFloat(s.borderRightWidth) || 0, c: s.borderRightColor, st: s.borderRightStyle },
      b: { w: parseFloat(s.borderBottomWidth) || 0, c: s.borderBottomColor, st: s.borderBottomStyle },
      l: { w: parseFloat(s.borderLeftWidth) || 0, c: s.borderLeftColor, st: s.borderLeftStyle },
    };
    const hasBg = bg && bg !== 'rgba(0, 0, 0, 0)' && bg !== 'transparent';
    const hasBorder = (border.t.w + border.r.w + border.b.w + border.l.w) > 0;

    if (hasBg || hasBorder) {
      shapes.push({
        x: rect.left - originX,
        y: rect.top - originY,
        w: rect.width,
        h: rect.height,
        bg: hasBg ? bg : null,
        borderRadius: parseFloat(s.borderTopLeftRadius) || 0,
        border: border,
        zIndex: parseInt(s.zIndex) || 0,
      });
    }

    for (const child of el.children) walkShapes(child);
  }

  // Inline tags whose text should be merged with their parent
  const INLINE_TAGS = new Set([
    'A','ABBR','B','BDI','BDO','BR','CITE','CODE','DATA','DFN','EM','I',
    'KBD','MARK','Q','RP','RT','RUBY','S','SAMP','SMALL','SPAN','STRONG',
    'SUB','SUP','TIME','U','VAR','WBR',
  ]);

  function isInline(el) {
    if (!(el instanceof Element)) return true; // text nodes are inline
    if (INLINE_TAGS.has(el.tagName)) return true;
    const d = getComputedStyle(el).display;
    return d === 'inline' || d === 'inline-block' || d === 'inline-flex';
  }

  function hasBlockChild(el) {
    for (const child of el.children) {
      if (!isInline(child)) return true;
    }
    return false;
  }

  function collectInlineText(el) {
    // Recursively collect text from inline children, preserving order
    let result = '';
    for (const node of el.childNodes) {
      if (node.nodeType === Node.TEXT_NODE) {
        result += node.nodeValue.replace(/\s+/g, ' ');
      } else if (node.nodeType === Node.ELEMENT_NODE && isInline(node)) {
        result += collectInlineText(node);
      }
    }
    return result;
  }

  function walkText(el) {
    if (!(el instanceof Element)) return;
    if (el.id === 'tweaks' || el.tagName === 'SCRIPT' || el.tagName === 'STYLE'
        || el.tagName === 'SVG' || el.tagName === 'svg') return;
    const s = getComputedStyle(el);
    if (s.display === 'none' || s.visibility === 'hidden') return;

    // If this element has block-level children, don't capture its own text —
    // recurse into children instead. This prevents parent elements from
    // duplicating text that their children will also emit.
    if (hasBlockChild(el)) {
      for (const child of el.children) walkText(child);
      return;
    }

    // This is a "leaf text container" — all children are inline.
    // Collect the full inline text content.
    const text = collectInlineText(el).trim();
    if (!text) return;

    const rect = el.getBoundingClientRect();
    if (!visible(el, rect)) return;

    texts.push({
      text,
      x: rect.left - originX,
      y: rect.top - originY,
      w: rect.width,
      h: rect.height,
      fontFamily: s.fontFamily,
      fontSize: parseFloat(s.fontSize),
      fontWeight: s.fontWeight,
      fontStyle: s.fontStyle,
      color: s.color,
      letterSpacing: s.letterSpacing,
      lineHeight: s.lineHeight,
      textAlign: s.textAlign,
      textTransform: s.textTransform,
    });
  }

  walkShapes(canvas);
  walkText(canvas);

  return {
    canvasWidth: canvas.offsetWidth,
    canvasHeight: canvas.offsetHeight,
    shapes,
    texts,
    speakerNotes,
  };
}
"""


# ----------------------------------------------------------------------
# Slide building
# ----------------------------------------------------------------------

@dataclass
class Options:
    width: int = 1920
    height: int = 1080
    wait_ms: int = 1500


def add_rect(slide, shape):
    """Draw a background rectangle with optional fill and border."""
    left = px_to_emu(shape["x"])
    top = px_to_emu(shape["y"])
    w = px_to_emu(max(shape["w"], 1))
    h = px_to_emu(max(shape["h"], 1))

    is_rounded = shape["borderRadius"] > 2
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if is_rounded else MSO_SHAPE.RECTANGLE
    pptx_shape = slide.shapes.add_shape(shape_type, left, top, w, h)

    # Fill
    bg = parse_rgb(shape["bg"]) if shape["bg"] else None
    if bg and bg[3] > 0.01:
        pptx_shape.fill.solid()
        # Approximate alpha by blending with white
        r, g, b, a = bg
        if a < 0.99:
            r = int(r * a + 255 * (1 - a))
            g = int(g * a + 255 * (1 - a))
            b = int(b * a + 255 * (1 - a))
        pptx_shape.fill.fore_color.rgb = RGBColor(r, g, b)
    else:
        pptx_shape.fill.background()

    # Border: suppress all shape outlines. The thin 1px CSS borders that
    # look subtle in a browser render as harsh outlines in PowerPoint.
    # Background fills alone are sufficient to convey the card/section
    # structure. This produces a much cleaner PPTX.
    pptx_shape.line.fill.background()

    return pptx_shape


def add_text(slide, t, opts: Options):
    left = px_to_emu(t["x"] - 1)
    top = px_to_emu(t["y"] - 1)
    w = px_to_emu(max(t["w"] + 4, 20))
    h = px_to_emu(max(t["h"] + 2, 12))

    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = css_align_to_pp(t["textAlign"])

    run = p.add_run()
    text = t["text"]
    if t.get("textTransform") == "uppercase":
        text = text.upper()
    elif t.get("textTransform") == "lowercase":
        text = text.lower()
    run.text = text

    font = run.font
    font.name = first_font(t["fontFamily"])
    # CSS px to PowerPoint pt: 1 px ≈ 0.75 pt
    font.size = Pt(t["fontSize"] * 0.75)

    try:
        w_val = int(t["fontWeight"])
    except (TypeError, ValueError):
        w_val = 700 if t.get("fontWeight") == "bold" else 400
    font.bold = w_val >= 600
    font.italic = (t.get("fontStyle") or "").startswith("italic")

    col = parse_rgb(t["color"])
    if col:
        font.color.rgb = rgb_to_color(col)

    return tb


def build_slide(prs: Presentation, data: dict, opts: Options):
    blank = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank)

    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Sort shapes: larger area first → smaller on top
    shapes = sorted(data["shapes"], key=lambda s: (s["w"] * s["h"]), reverse=True)
    for s in shapes:
        try:
            add_rect(slide, s)
        except Exception as e:
            print(f"  ! shape skipped: {e}", file=sys.stderr)

    for t in data["texts"]:
        try:
            add_text(slide, t, opts)
        except Exception as e:
            print(f"  ! text skipped: {e}", file=sys.stderr)

    # Speaker notes
    notes = data.get("speakerNotes", "").strip()
    if notes:
        try:
            slide.notes_slide.notes_text_frame.text = notes
        except Exception:
            pass


# ----------------------------------------------------------------------
# Orchestration
# ----------------------------------------------------------------------

async def capture_slide(browser: Browser, path: Path, opts: Options) -> dict:
    context = await browser.new_context(
        viewport={"width": opts.width, "height": opts.height},
        device_scale_factor=1,
    )
    page = await context.new_page()
    url = path.resolve().as_uri()
    await page.goto(url, wait_until="networkidle")
    # Wait for web fonts
    await page.evaluate("() => document.fonts ? document.fonts.ready : Promise.resolve()")
    await page.wait_for_timeout(opts.wait_ms)
    data = await page.evaluate(EXTRACT_JS)
    await context.close()
    return data


async def run(inputs: list[Path], out_path: Path, opts: Options):
    prs = Presentation()
    prs.slide_width = px_to_emu(opts.width)
    prs.slide_height = px_to_emu(opts.height)

    async with async_playwright() as pw:
        browser = await pw.chromium.launch(headless=True)
        for i, path in enumerate(inputs, 1):
            print(f"[{i}/{len(inputs)}] {path.name}")
            data = await capture_slide(browser, path, opts)
            print(f"    {len(data['shapes'])} shapes, {len(data['texts'])} text runs")
            build_slide(prs, data, opts)
        await browser.close()

    prs.save(str(out_path))
    print(f"\nwrote {out_path}  ({len(inputs)} slide{'s' if len(inputs) != 1 else ''})")


def expand_inputs(patterns: list[str]) -> list[Path]:
    out: list[Path] = []
    for p in patterns:
        matches = sorted(glob.glob(p))
        if not matches and Path(p).exists():
            matches = [p]
        for m in matches:
            out.append(Path(m))
    if not out:
        print(f"no files matched: {patterns}", file=sys.stderr)
        sys.exit(2)
    return out


def main():
    ap = argparse.ArgumentParser(description="Convert HTML slides to an editable PPTX.")
    ap.add_argument("inputs", nargs="+", help="HTML files or globs (e.g. 'slides/*.html')")
    ap.add_argument("-o", "--output", default="deck.pptx", help="output .pptx path")
    ap.add_argument("--width", type=int, default=1920, help="slide width in px (default 1920)")
    ap.add_argument("--height", type=int, default=1080, help="slide height in px (default 1080)")
    ap.add_argument("--wait", type=int, default=1500, help="ms to wait after load before extracting (default 1500)")
    args = ap.parse_args()

    files = expand_inputs(args.inputs)
    opts = Options(width=args.width, height=args.height, wait_ms=args.wait)
    asyncio.run(run(files, Path(args.output), opts))


if __name__ == "__main__":
    main()
